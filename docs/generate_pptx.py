"""Generate GATRA World Monitor User Guide PPTX presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ── Colors ──────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Deep navy background
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)       # Card background
ACCENT_GREEN = RGBColor(0x22, 0xC5, 0x5E)  # Status green
ACCENT_ORANGE = RGBColor(0xF9, 0x73, 0x16) # Warning orange
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)    # Critical red
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)   # Info blue
ACCENT_YELLOW = RGBColor(0xEA, 0xB3, 0x08) # Medium yellow
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6) # A2A purple
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)
TEXT_HEADING = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=TEXT_DIM, bullet_color=ACCENT_GREEN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_card(slide, left, top, width, height, title, body_lines,
             accent=ACCENT_GREEN):
    """Add a rounded-rectangle card with title + bullet body."""
    shape = slide.shapes.add_shape(
        5, left, top, width, height  # 5 = rounded rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(12)
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)

    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = accent
    p.font.name = "Calibri"
    p.space_after = Pt(8)

    # Body lines
    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DIM
        p.font.name = "Calibri"
        p.space_after = Pt(3)

    return shape


def section_header(slide, number, title, subtitle=""):
    """Add a section header with number badge."""
    set_slide_bg(slide)
    # Section number circle
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(1.2), Inches(1.2),
                str(number), font_size=48, color=ACCENT_GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, Inches(2.2), Inches(2.4), Inches(9), Inches(1),
                title, font_size=40, color=TEXT_WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(2.2), Inches(3.5), Inches(9), Inches(1),
                    subtitle, font_size=18, color=TEXT_DIM)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "GATRA World Monitor", font_size=48, color=TEXT_WHITE, bold=True)
add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
            "User Guide & Navigation Reference", font_size=28, color=ACCENT_GREEN)
add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
            "https://worldmonitor-gatra.vercel.app", font_size=16, color=ACCENT_BLUE)
add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
            "Cyber Variant (GATRA SOC Edition)  |  March 2026", font_size=14, color=TEXT_DIM)
add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
            "AI-Driven Security Operations Center for Critical Infrastructure",
            font_size=14, color=TEXT_DIM)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2: Table of Contents
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
            "Table of Contents", font_size=32, color=TEXT_WHITE, bold=True)

left_items = [
    "1.  Getting Started",
    "2.  Dashboard Overview",
    "3.  Interactive Map & Layers",
    "4.  GATRA SOC Panel",
    "5.  SOC COMMS (AI Agent Chat)",
    "6.  A2A Security Monitor",
    "7.  Social Threat Intel",
    "8.  CVE Feed",
]
right_items = [
    "9.   IoC Lookup",
    "10.  Ransomware Tracker",
    "11.  CII Monitor",
    "12.  Prediction Signals",
    "13.  Security Posture",
    "14.  Customization & Settings",
    "15.  URL Deep Linking",
    "16.  Data Sources & Refresh Rates",
]
add_bullet_list(slide, Inches(1), Inches(1.6), Inches(5.5), Inches(5),
                left_items, font_size=16, color=TEXT_HEADING)
add_bullet_list(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5),
                right_items, font_size=16, color=TEXT_HEADING)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3: Getting Started
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 1, "Getting Started",
               "First-time setup and system requirements")

# ════════════════════════════════════════════════════════════════════
# SLIDE 4: Getting Started — Details
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Getting Started", font_size=28, color=TEXT_WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5),
         "First-Time Setup", [
             "1. Open worldmonitor-gatra.vercel.app",
             "2. Wait 5-15s for live data to load",
             "3. Check sidebar for available panels",
             "4. Adjust time range (1h / 6h / 24h / 7d)",
             "5. Toggle map layers via layer control",
             "6. Try SOC Chat: type a cyber question",
         ], ACCENT_GREEN)

add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2.5),
         "System Requirements", [
             "Recommended: Desktop, 1280px+ screen",
             "Best browsers: Chrome 120+, Firefox 120+",
             "Tablet: Supported (simplified layout)",
             "Mobile: Basic support (best on desktop)",
             "Dark mode is the default theme",
         ], ACCENT_BLUE)

add_card(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5),
         "Data Freshness Indicators", [
             'LIVE (green pulsing) = Real-time data from API',
             'MOCK (yellow) = Simulated data (API temporarily unavailable)',
             'Each panel header shows item count + data source badge',
             'Hover (i) icon on any panel for description + refresh interval',
         ], ACCENT_YELLOW)

# ════════════════════════════════════════════════════════════════════
# SLIDE 5: Dashboard Overview
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 2, "Dashboard Overview",
               "Layout, panels, and navigation")

# ════════════════════════════════════════════════════════════════════
# SLIDE 6: Dashboard Layout
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Dashboard Layout", font_size=28, color=TEXT_WHITE, bold=True)

# ASCII layout mockup
add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(4),
            "+-------------------------------------------------+\n"
            "|  HEADER BAR (title, settings, theme, language)   |\n"
            "+--------+----------------------------------------+\n"
            "|        |                                        |\n"
            "| PANELS |         INTERACTIVE MAP                |\n"
            "| (left  |    (global view with layers)           |\n"
            "|  side, |                                        |\n"
            "| scroll)|   Click markers for details            |\n"
            "|        |   Layer toggles top-right              |\n"
            "|        |   Zoom +/- controls                    |\n"
            "+--------+----------------------------------------+",
            font_size=14, color=ACCENT_GREEN,
            font_name="Courier New")

add_bullet_list(slide, Inches(1), Inches(5.5), Inches(11), Inches(1.5), [
    "Panels: scrollable column on the left side, each collapsible",
    "Map: interactive global view in the center, supports zoom/pan/click",
    "Header: title, language selector, theme toggle, settings",
], font_size=13)

# ════════════════════════════════════════════════════════════════════
# SLIDE 7: Interactive Map
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 3, "Interactive Map & Layers",
               "Global visualization of cybersecurity events and infrastructure")

# ════════════════════════════════════════════════════════════════════
# SLIDE 8: Map Layers Table
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Map Layers (Cyber Variant)", font_size=28, color=TEXT_WHITE, bold=True)

layers_on = [
    ("GATRA Alerts", "SOC alert geolocations by severity"),
    ("Cyber Threats", "Malware hosts, C2 servers, phishing URLs"),
    ("Conflicts", "Armed conflict zones (UCDP data)"),
    ("Undersea Cables", "Communication cable infrastructure"),
    ("Outages", "Internet service disruptions"),
    ("Datacenters", "AI/GPU compute facilities worldwide"),
    ("Military", "Aircraft & vessel tracking"),
    ("Natural Disasters", "Earthquakes, volcanoes, storms"),
    ("Fires", "Satellite-detected wildfires"),
]

for i, (name, desc) in enumerate(layers_on):
    y = Inches(1.3) + Inches(i * 0.6)
    add_textbox(slide, Inches(1), y, Inches(0.6), Inches(0.4),
                "ON", font_size=10, color=BG_DARK, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Green dot
    dot = slide.shapes.add_shape(9, Inches(1.1), y + Pt(4), Pt(12), Pt(12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_GREEN
    dot.line.fill.background()

    add_textbox(slide, Inches(1.8), y, Inches(3), Inches(0.4),
                name, font_size=14, color=TEXT_WHITE, bold=True)
    add_textbox(slide, Inches(5), y, Inches(7.5), Inches(0.4),
                desc, font_size=12, color=TEXT_DIM)

add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
            "Additional layers available: Flights, Pipelines, Nuclear, Bases, Hotspots, Weather, and more",
            font_size=12, color=TEXT_DIM)

# ════════════════════════════════════════════════════════════════════
# SLIDE 9: GATRA SOC Panel
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 4, "GATRA SOC Panel",
               "AI-driven alert triage with Asset Relevance Scoring")

# ════════════════════════════════════════════════════════════════════
# SLIDE 10: GATRA SOC — Agent Status
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "GATRA SOC — 5-Agent Pipeline", font_size=28, color=TEXT_WHITE, bold=True)

agents = [
    ("ADA", "Anomaly Detection Agent", "Isolation Forest + LSTM ensemble", ACCENT_GREEN),
    ("TAA", "Threat Analysis Agent", "Actor-Critic RL triage & prioritization", ACCENT_ORANGE),
    ("CRA", "Containment & Response", "Automated playbooks & containment", ACCENT_RED),
    ("CLA", "Compliance & Logging", "Audit trails & regulatory compliance", ACCENT_BLUE),
    ("RVA", "Risk & Vulnerability", "Vulnerability assessment & exposure", ACCENT_PURPLE),
]

for i, (code, name, desc, color) in enumerate(agents):
    x = Inches(0.5) + Inches(i * 2.5)
    add_card(slide, x, Inches(1.4), Inches(2.3), Inches(2.2),
             code, [name, "", desc], color)

add_card(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(2.8),
         "Stats Row", [
             "Active: Number of critical + high severity incidents",
             "MTTR: Mean Time to Respond (green <15m, yellow <30m, red >30m)",
             "24h Alerts: Total alerts in selected time range",
             "24h Resp: CRA automated response actions taken",
             "",
             "Time Range Filter: 1h | 6h | 24h | 48h | 7d | all",
         ], ACCENT_YELLOW)

# ════════════════════════════════════════════════════════════════════
# SLIDE 11: GATRA SOC — Relevance Scoring
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Asset Relevance Scoring", font_size=28, color=TEXT_WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
            "Transforms 50 generic alerts into prioritized, actionable intelligence",
            font_size=16, color=TEXT_DIM)

# Scoring breakdown
add_card(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3),
         "Scoring Breakdown (max 100 pts)", [
             "Vendor match:         40 pts  (uses your tech stack)",
             "Product match:        30 pts  (specific product hit)",
             "Industry keyword:     15 pts  (telecom, 5G, etc.)",
             "Ransomware campaign:  10 pts  (known ransomware use)",
             "Due date urgency:      5 pts  (CISA deadline <7 days)",
             "",
             'Example: Cisco ASA vuln + ransomware = 40+30+10 = 80 "MATCH"',
         ], ACCENT_GREEN)

# Badge meanings
add_card(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(3),
         "Relevance Badges", [
             "MATCH (green, 80-100)",
             "  Direct hit: vendor AND product in your stack",
             "",
             "RELATED (yellow, 50-79)",
             "  Partial: vendor match, no specific product",
             "",
             "LOW (gray, 20-49) / INFO (dark, 0-19)",
             "  General advisory, not in your stack",
         ], ACCENT_YELLOW)

# Features
add_card(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2),
         "Panel Features", [
             '"X of Y alerts affect your infrastructure" summary bar with progress indicator',
             "Green vendor badges showing matched vendors (Microsoft, Cisco, Fortinet, etc.)",
             '"Show Relevant Only" toggle button to filter to MATCH + RELATED alerts only',
             "Alerts sorted by relevance first, then severity (most relevant at top)",
             "Each alert shows matched assets: e.g., Cisco -> SD-WAN",
         ], ACCENT_BLUE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 12: SOC COMMS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 5, "SOC COMMS — AI Agent Chat",
               "Interactive chat with 5 specialized cybersecurity AI agents")

# ════════════════════════════════════════════════════════════════════
# SLIDE 13: SOC COMMS — How to Use
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Talking to GATRA Agents", font_size=28, color=TEXT_WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.8),
         "Ask Natural Questions", [
             '"any new anomaly?" -> ADA responds',
             '"any real threats?" -> TAA responds',
             '"block 10.0.0.1" -> CRA responds',
             '"generate report" -> CLA responds',
             '"CVE status?" -> RVA responds',
             "",
             "Agents auto-detect keywords in your message",
         ], ACCENT_GREEN)

add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2.8),
         "Direct Commands", [
             "@ADA — Target Anomaly Detection Agent",
             "@TAA — Target Threat Analysis Agent",
             "@CRA — Target Containment Response Agent",
             "@CLA — Target Compliance & Logging Agent",
             "@RVA — Target Risk & Vulnerability Agent",
             "",
             "/help — Show all available commands",
         ], ACCENT_BLUE)

add_card(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5),
         "IOC Lookups in Chat", [
             "Paste any indicator directly into chat for instant threat intelligence:",
             "",
             "  IP address:  192.168.1.1  ->  Queries ThreatFox + AbuseIPDB",
             "  Domain:      evil-site.com  ->  Queries URLhaus + ThreatFox",
             "  Hash:        d41d8cd98f...  ->  Queries MalwareBazaar",
             "  URL:         https://phishing.com/login  ->  Queries URLhaus",
         ], ACCENT_ORANGE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 14: A2A Security Monitor
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 6, "A2A Security Monitor",
               "Agent-to-Agent trust scoring, traffic monitoring, and validation")

# ════════════════════════════════════════════════════════════════════
# SLIDE 15: A2A Details
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "A2A Security Features", font_size=28, color=TEXT_WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(1.3), Inches(3.7), Inches(2.5),
         "Agent Registry", [
             "All registered agents listed",
             "Online/offline status dots",
             "Trust score bars (0-100)",
             "",
             "85-100: Trusted (internal)",
             "60-84: Verified (external)",
             "40-59: Pending review",
             "0-39: Blocked/malicious",
         ], ACCENT_PURPLE)

add_card(slide, Inches(5), Inches(1.3), Inches(3.7), Inches(2.5),
         "Traffic Feed", [
             "Live A2A request stream:",
             "",
             "Clean (green) = Verified",
             "Suspicious (yellow) = Drift",
             "Blocked (red) = Rejected",
             "",
             "Shows: time, route,",
             "skill/method, verdict",
         ], ACCENT_ORANGE)

add_card(slide, Inches(9.2), Inches(1.3), Inches(3.3), Inches(2.5),
         "Security Tests", [
             "Interactive test grid:",
             "- Certificate validation",
             "- Encryption check",
             "- Injection resistance",
             "- Protocol compliance",
             "",
             "Click Run to execute",
             "Pass/fail/warn results",
         ], ACCENT_GREEN)

add_card(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5),
         "CII-Aware Trust Policy", [
             "Trust thresholds automatically adjust based on Country Instability Index:",
             "",
             "Standard (CII normal):   Relaxed trust, allow unsigned, 100 req/hr",
             "Elevated (CII rising):   Min trust 60, reject unsigned, manual approval, 30 req/hr",
             "Critical (CII spike):    Min trust 85, block low-trust regions, 10 req/hr",
         ], ACCENT_RED)

# ════════════════════════════════════════════════════════════════════
# SLIDE 16: Social Threat Intel
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 7, "Social Threat Intel",
               "Real-time cybersecurity chatter from Hacker News, Mastodon, and Bluesky")

# ════════════════════════════════════════════════════════════════════
# SLIDE 17: Social Threat Details
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Social Threat Intel Features", font_size=28, color=TEXT_WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(1.3), Inches(3.7), Inches(2.5),
         "Hacker News (orange)", [
             "Top cybersecurity stories",
             "Source: Algolia search API",
             "Keywords: vulnerability,",
             "  breach, exploit, security",
             "Free, no auth required",
         ], ACCENT_ORANGE)

add_card(slide, Inches(5), Inches(1.3), Inches(3.7), Inches(2.5),
         "Mastodon (purple)", [
             "#cybersecurity + #infosec",
             "Sources: mastodon.social",
             "  and hachyderm.io",
             "Public hashtag timelines",
             "Free, no auth required",
         ], ACCENT_PURPLE)

add_card(slide, Inches(9.2), Inches(1.3), Inches(3.3), Inches(2.5),
         "Bluesky (blue)", [
             "Cybersecurity keyword",
             "  search via AT Protocol",
             "public.api.bsky.app",
             "Free, no auth required",
             "(May be WAF-limited)",
         ], ACCENT_BLUE)

add_card(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8),
         "Panel Controls & Features", [
             "Platform filter tabs: All | Bluesky | HN | Mastodon (click to filter)",
             "Stats row: Posts count | Threats detected | Total engagement | Trending keyword",
             "Threat keywords auto-highlighted in red: CVE, ransomware, zero-day, breach, exploit, malware",
             "Each post shows: author, content snippet, engagement (likes/reposts/replies), timestamp",
             "Click any post to open the original on its platform",
             "Refresh: every 10 minutes | Up to 50 posts displayed",
         ], ACCENT_GREEN)

# ════════════════════════════════════════════════════════════════════
# SLIDE 18: CVE + IoC + Ransomware (combined)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Threat Intelligence Panels", font_size=28, color=TEXT_WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(1.3), Inches(3.7), Inches(3),
         "8. CVE Feed", [
             "Source: NVD v2 + CISA KEV",
             "Refresh: every 10 minutes",
             "",
             "CVSS score badges (colored)",
             "KEV = actively exploited",
             "Affected products listed",
             "Stats: total, critical, high,",
             "  exploited counts",
         ], ACCENT_RED)

add_card(slide, Inches(5), Inches(1.3), Inches(3.7), Inches(3),
         "9. IoC Lookup", [
             "Search: IP, domain, hash, URL",
             "Auto-detects indicator type",
             "Sources: ThreatFox, URLhaus,",
             "  MalwareBazaar",
             "",
             "Returns: threat level, verdicts,",
             "  malware family, timeline,",
             "  related IOCs (click to pivot)",
         ], ACCENT_BLUE)

add_card(slide, Inches(9.2), Inches(1.3), Inches(3.3), Inches(3),
         "10. Ransomware Tracker", [
             "Source: ransomware.live",
             "Refresh: every 5 minutes",
             "",
             "30-day victim count",
             "Top groups bar chart",
             "Recent victims feed:",
             "  group, name, country,",
             "  sector, date",
         ], ACCENT_ORANGE)

add_card(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.3),
         "IoC Lookup — How to Search", [
             "1. Type or paste an indicator into the search box",
             "2. Auto-detected type badge appears: IP (blue) | Domain (purple) | Hash (pink) | URL (amber)",
             "3. Click Search or press Enter",
             "4. Results show: Threat Level (Malicious/Suspicious/Clean) + source verdicts + tags + timeline",
             "5. Click related IOCs to pivot search and investigate further",
         ], ACCENT_GREEN)

# ════════════════════════════════════════════════════════════════════
# SLIDE 19: CII + Prediction Signals
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Geopolitical Intelligence Panels", font_size=28, color=TEXT_WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3),
         "11. CII Monitor (Country Instability Index)", [
             "Indonesia focus card: score (0-100), trend arrow, sparkline",
             "Color: Green (<30) | Yellow (30-59) | Orange (60-79) | Red (80+)",
             "",
             "R_geo Impact Signal: how instability feeds RL reward function",
             "  Active (yellow) = elevated, Nominal (green) = normal",
             "",
             "Regional grid: SG, MY, PH, TH, VN, AU, PNG, MM, CN",
             "Each shows: flag, score, trend delta, mini progress bar",
         ], ACCENT_YELLOW)

add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(3),
         "12. Prediction Signals", [
             "Source: Prediction markets (Polymarket)",
             "",
             "Early Warning Multiplier (1.0 - 3.0):",
             "  1.0 = nominal, 2.0+ = elevated risk",
             "",
             "Per-market cards show:",
             "  Probability bar, threat level, relevance tier",
             "  Tier 1: Indonesia direct | Tier 2: ASEAN",
             "  Tier 3: Global conflict | Tier 4: Economic",
             "  Volume, velocity (rate of change)",
         ], ACCENT_PURPLE)

add_card(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.3),
         "13. Security Posture (Personal)", [
             "Self-assessment of analyst cybersecurity hygiene",
             "Categories: Authentication (passwords, 2FA) | Device Security (OS, antivirus, encryption) | Online Behavior (VPN, updates, backups)",
             "Auto-scan: browser version, HTTPS connection, email domain security",
             "Score: 0-100 with color grade (Strong 80+ / Moderate 60-79 / Weak <60)",
         ], ACCENT_GREEN)

# ════════════════════════════════════════════════════════════════════
# SLIDE 20: Customization & URL Deep Linking
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Customization & URL Deep Linking", font_size=28, color=TEXT_WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5),
         "14. Settings & Customization", [
             "Theme: Dark / Light mode toggle",
             "Language: EN, ID, ES, FR, DE, JA, AR...",
             "Panel visibility: sidebar toggles",
             "Map layers: layer icon on map",
             "Panel order: drag to reorder",
             "All preferences saved to browser",
         ], ACCENT_BLUE)

add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2.5),
         "15. URL Parameters", [
             "?view=global  (or us)",
             "&zoom=4  (1-20)",
             "&center=10.5,106.8  (lat,lon)",
             "&layers=gatraAlerts,cyberThreats",
             "&timeRange=24h  (1h/6h/24h/7d/all)",
             "&panels=gatra-soc,cve-feed",
             "&theme=dark  (dark/light)",
         ], ACCENT_GREEN)

add_card(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8),
         "Useful Bookmarks", [
             "Full Cyber Dashboard:  ?view=global&zoom=2.5&layers=gatraAlerts,cyberThreats,cables,outages",
             "Indonesia Focus:       ?view=global&zoom=5&center=-2.5,118&layers=gatraAlerts,cyberThreats&country=ID",
             "Gulf Region:           ?view=global&zoom=5&center=25,52&layers=flights,cables,conflicts",
             "SOC Analyst View:      ?panels=gatra-soc,ioc-lookup,ransomware-tracker&timeRange=24h",
         ], ACCENT_ORANGE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 21: Data Sources & Refresh Rates
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "16. Data Sources & Refresh Rates", font_size=28, color=TEXT_WHITE, bold=True)

sources = [
    ("GATRA SOC", "CISA KEV catalog", "60 seconds"),
    ("CVE Feed", "NVD v2 API", "10 minutes"),
    ("Social Threat Intel", "HN + Mastodon + Bluesky", "10 minutes"),
    ("IoC Lookup", "ThreatFox / URLhaus / MalwareBazaar", "On-demand (5m cache)"),
    ("Ransomware Tracker", "ransomware.live API", "5 minutes"),
    ("CII Monitor", "Prediction markets + geo feeds", "5 minutes"),
    ("Prediction Signals", "Polymarket / prediction APIs", "5 minutes"),
    ("A2A Security", "Internal agent traffic", "Real-time"),
    ("Map Layers", "AIS, flight, UCDP, USGS, etc.", "2-10 minutes"),
]

# Header
add_textbox(slide, Inches(1), Inches(1.3), Inches(3.5), Inches(0.4),
            "Panel", font_size=12, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(4.5), Inches(1.3), Inches(5), Inches(0.4),
            "Data Source", font_size=12, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(10), Inches(1.3), Inches(2.5), Inches(0.4),
            "Refresh", font_size=12, color=ACCENT_GREEN, bold=True)

for i, (panel, source, refresh) in enumerate(sources):
    y = Inches(1.8) + Inches(i * 0.55)
    color = TEXT_WHITE if i % 2 == 0 else TEXT_DIM
    add_textbox(slide, Inches(1), y, Inches(3.5), Inches(0.4),
                panel, font_size=13, color=color, bold=(i % 2 == 0))
    add_textbox(slide, Inches(4.5), y, Inches(5), Inches(0.4),
                source, font_size=12, color=TEXT_DIM)
    add_textbox(slide, Inches(10), y, Inches(2.5), Inches(0.4),
                refresh, font_size=12, color=ACCENT_BLUE)

add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
            "Background tabs automatically reduce polling by 4x to save bandwidth",
            font_size=12, color=TEXT_DIM)

# ════════════════════════════════════════════════════════════════════
# SLIDE 22: Closing
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1),
            "GATRA World Monitor", font_size=44, color=TEXT_WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
            "AI-Driven Security Operations Center\nfor Critical Infrastructure",
            font_size=22, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
            "https://worldmonitor-gatra.vercel.app",
            font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
            "For support and feedback, contact the GATRA development team",
            font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────
out = Path(__file__).parent / "GATRA_WorldMonitor_User_Guide.pptx"
prs.save(str(out))
print(f"Saved: {out} ({out.stat().st_size / 1024:.0f} KB, {len(prs.slides)} slides)")
